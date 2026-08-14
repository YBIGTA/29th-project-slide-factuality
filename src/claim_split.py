#!/usr/bin/env python3
"""PPTX의 텍스트를 빠르게 추출해 claim JSONL로 저장한다.

추출 기준은 ``prompts/claim_split_v1.txt``에 문서화되어 있다.
"""

from __future__ import annotations

import argparse
import json
import re
import statistics
import sys
from dataclasses import asdict, dataclass, field, replace
from pathlib import Path
from typing import Any, Iterable, Sequence

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


페이지_표시 = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")
출처_표시 = re.compile(
    r"(?:출처|source|references?|doi)\s*[:：]|한국예술연구|Korean\s+J\.|\|\s*제?\d+호",
    re.IGNORECASE,
)
목록_표시 = re.compile(r"^\s*(?:[✓✔✗✕•·▪◦●○]|[-–—]|\d+[.)](?=\s))\s*")
목차_제목 = re.compile(r"^(?:발표\s*구성|목차|차례)$", re.IGNORECASE)


@dataclass(frozen=True)
class 텍스트항목:
    """도형에서 뽑아낸 텍스트 한 덩어리. item_id/role/tag는 나중 단계에서 채운다."""

    text: str
    top: int
    left: int
    width: int
    height: int
    font_size: float
    source_type: str = "text"
    item_id: int = 0
    role: str = "body"
    tag: str = "본문"


@dataclass(frozen=True)
class 슬라이드정보:
    number: int
    title: str
    items: list[텍스트항목]
    warnings: list[str] = field(default_factory=list)


def clean_text(text: str) -> str:
    """불필요한 공백은 줄이고 PPT의 줄바꿈은 보존한다."""
    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def max_font_size(shape: Any) -> float:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return 0.0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    return max(sizes, default=0.0)


def shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return clean_text("\n".join(p.text for p in shape.text_frame.paragraphs))


def table_items(shape: Any) -> list[텍스트항목]:
    """표는 셀 단위가 아니라 행 단위 텍스트로 읽는다."""
    rows = list(shape.table.rows)
    row_height = int(shape.height / max(len(rows), 1))
    result: list[텍스트항목] = []
    for index, row in enumerate(rows):
        cells = [clean_text(cell.text) for cell in row.cells]
        text = " | ".join(cell for cell in cells if cell)
        if text:
            result.append(
                텍스트항목(
                    text,
                    int(shape.top) + index * row_height,
                    int(shape.left),
                    int(shape.width),
                    row_height,
                    max_font_size(row.cells[0]) if row.cells else 0.0,
                    "table_row",
                )
            )
    return result


def chart_items(shape: Any, warnings: list[str]) -> list[텍스트항목]:
    """네이티브 차트의 제목과 범주별 값을 텍스트로 바꾼다."""
    result: list[텍스트항목] = []
    chart = shape.chart
    if getattr(chart, "has_title", False):
        title = clean_text(chart.chart_title.text_frame.text)
        if title:
            result.append(
                텍스트항목(
                    title,
                    int(shape.top),
                    int(shape.left),
                    int(shape.width),
                    int(shape.height * 0.15),
                    0.0,
                    "chart_title",
                )
            )
    try:
        for plot_index, plot in enumerate(chart.plots, start=1):
            categories = [clean_text(str(value)) for value in plot.categories]
            series_list = list(plot.series)
            for category_index, category in enumerate(categories):
                values: list[str] = []
                for series in series_list:
                    series_values = list(series.values)
                    if category_index < len(series_values):
                        name = clean_text(str(series.name or "계열"))
                        values.append(f"{name}: {series_values[category_index]}")
                result.append(
                    텍스트항목(
                        " | ".join([category, *values]),
                        int(shape.top + shape.height * (category_index + 1) / (len(categories) + 1)),
                        int(shape.left),
                        int(shape.width),
                        int(shape.height / max(len(categories), 1)),
                        0.0,
                        f"chart_data_{plot_index}",
                    )
                )
    except (AttributeError, TypeError, ValueError) as exc:
        warnings.append(f"chart_data_not_extracted:{type(exc).__name__}")
    return result


def shape_items(shape: Any, warnings: list[str]) -> list[텍스트항목]:
    """그룹·표·차트·텍스트 도형을 같은 항목 형식으로 변환한다."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        result: list[텍스트항목] = []
        for child in shape.shapes:
            result.extend(shape_items(child, warnings))
        return result
    if getattr(shape, "has_table", False):
        return table_items(shape)
    if getattr(shape, "has_chart", False):
        return chart_items(shape, warnings)

    text = shape_text(shape)
    if text:
        return [
            텍스트항목(
                text,
                int(shape.top),
                int(shape.left),
                int(shape.width),
                int(shape.height),
                max_font_size(shape),
            )
        ]
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        warnings.append("image_without_machine_readable_text")
    return []


def is_page_number(item: 텍스트항목, slide_width: int, slide_height: int) -> bool:
    if 페이지_표시.fullmatch(item.text):
        return True
    return bool(
        re.fullmatch(r"\d{1,3}", item.text)
        and item.top >= slide_height * 0.85
        and item.left >= slide_width * 0.75
    )


def find_title(items: Sequence[텍스트항목], slide_height: int) -> int | None:
    candidates = [
        item
        for item in items
        if item.source_type in {"text", "chart_title"}
        and item.top < slide_height * 0.40
        and not re.fullmatch(r"\d{1,4}(?:[–-]\d{1,4})?", item.text)
    ]
    if not candidates:
        return None
    return max(candidates, key=lambda item: (item.font_size, -item.top)).item_id


def default_tag(item: 텍스트항목, title: 텍스트항목 | None, median_font: float, slide_height: int, is_toc: bool) -> str:
    if item.role == "title":
        return "제목"
    if 출처_표시.search(item.text):
        return "출처표기"
    if is_toc:
        return "목차항목"
    if item.source_type == "table_row" or item.source_type.startswith("chart_data_"):
        return "행"
    if item.source_type == "chart_title":
        return "카드헤더"
    if re.match(r"^(?:18|19|20)\d{2}(?:[–—-]\d{2,4})?\s*\S", item.text) and len(item.text) > 6:
        return "연표행"
    if re.fullmatch(r"[+−-]?\d+(?:[.,]\d+)?\s*(?:%|배|명|개|건|점|종|균주)?", item.text):
        return "스탯"
    if 목록_표시.match(item.text):
        return "불릿"
    if re.match(r'^[“”"‘’\']', item.text):
        return "인용"
    if title is not None:
        gap = item.top - (title.top + title.height)
        aligned = (
            abs(item.left - title.left) <= max(title.width * 0.15, 100_000)
            and item.width >= title.width * 0.55
        )
        large = title.font_size > 0 and item.font_size >= title.font_size * 0.35
        if 0 <= gap <= slide_height * 0.14 and aligned and large and len(item.text) <= 120:
            return "부제"
    if median_font > 0 and item.font_size >= median_font * 1.2 and len(item.text) <= 45:
        return "카드헤더"
    if len(item.text) >= 55 or re.search(r"[.!?。]$", item.text):
        return "카드본문"
    return "본문"


def attach_tags(items: Sequence[텍스트항목], title_id: int | None, slide_height: int, title_text: str) -> list[텍스트항목]:
    title = next((item for item in items if item.item_id == title_id), None)
    fonts = [item.font_size for item in items if item.item_id != title_id and item.font_size > 0]
    median_font = statistics.median(fonts) if fonts else 0.0
    is_toc = bool(목차_제목.fullmatch(title_text))

    def group_key(item: 텍스트항목) -> tuple[int, int, int, float]:
        return (item.left, item.width, item.height, round(item.font_size, 1))

    # 같은 크기와 위치로 반복된 도형은 카드 헤더 또는 불릿일 가능성이 높다.
    groups: dict[tuple[int, int, int, float], int] = {}
    for item in items:
        groups[group_key(item)] = groups.get(group_key(item), 0) + 1

    result: list[텍스트항목] = []
    for item in items:
        current = replace(item, role="title" if item.item_id == title_id else "body")
        tag = default_tag(current, title, median_font, slide_height, is_toc)
        if tag == "본문" and groups.get(group_key(item), 0) >= 2:
            tag = "카드헤더" if item.font_size >= median_font * 1.05 and len(item.text) <= 35 else "불릿"
        result.append(replace(current, tag=tag))
    return result


def read_slides(path: Path) -> list[슬라이드정보]:
    presentation = Presentation(path)
    width, height = int(presentation.slide_width), int(presentation.slide_height)
    result: list[슬라이드정보] = []

    for number, slide in enumerate(presentation.slides, start=1):
        warnings: list[str] = []
        raw: list[텍스트항목] = []
        for shape in slide.shapes:
            raw.extend(shape_items(shape, warnings))
        raw = [item for item in raw if item.text and not is_page_number(item, width, height)]
        raw.sort(key=lambda item: (item.top, item.left, item.source_type))
        items = [replace(item, item_id=index) for index, item in enumerate(raw, start=1)]
        title_id = find_title(items, height)
        title = next((item.text for item in items if item.item_id == title_id), "")
        tagged = attach_tags(items, title_id, height, title)
        result.append(슬라이드정보(number, title, tagged, sorted(set(warnings))))
    return result


def split_claims(text: str) -> list[str]:
    """일반 상자는 유지하고, 여러 불릿이 든 상자만 항목별로 나눈다."""
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(lines) > 1 and all(목록_표시.match(line) for line in lines):
        return [목록_표시.sub("", line).strip() for line in lines]
    return [목록_표시.sub("", text).strip()] if text.strip() else []


def build_records(slide: 슬라이드정보) -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    seen: set[str] = set()
    # Context는 현재 claim과 제목을 포함한 슬라이드 전체 텍스다.
    # 따라서 같은 슬라이드에서 추출된 모든 claim은 동일한 Context를
    # 가지며, 항목 순서를 통해 claim의 상대적 위치도 확인할 수 있다.
    context = [f"[{item.tag}] {item.text}" for item in slide.items]
    for source in slide.items:
        for claim in split_claims(source.text):
            if not claim or claim in seen:
                continue
            seen.add(claim)
            result.append(
                {
                    "Slide #": slide.number,
                    "Claim (PPT)": claim,
                    "Slide_Title": slide.title,
                    "Context (PPT)": list(context),
                }
            )
    return result


def save_jsonl(path: Path, records: Iterable[dict[str, Any]]) -> int:
    rows = list(records)
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix(path.suffix + ".tmp")
    with temporary.open("w", encoding="utf-8", newline="\n") as stream:
        for row in rows:
            stream.write(json.dumps(row, ensure_ascii=False) + "\n")
    temporary.replace(path)
    return len(rows)


def find_inputs(inputs: Sequence[Path]) -> list[Path]:
    def is_deck(path: Path) -> bool:
        return path.is_file() and path.suffix.lower() == ".pptx" and not path.name.startswith("~$")

    found: list[Path] = []
    for source in inputs:
        if source.is_dir():
            found.extend(path for path in source.iterdir() if is_deck(path))
        elif is_deck(source):
            found.append(source)
        else:
            raise FileNotFoundError(f"PPTX를 찾을 수 없습니다: {source}")
    unique = {path.resolve(): path for path in found}
    return [unique[key] for key in sorted(unique, key=lambda path: str(path).casefold())]


def parse_arguments(argv: Sequence[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("inputs", type=Path, nargs="*", default=[Path("decks")], help="PPTX 파일 또는 폴더")
    parser.add_argument("--output", type=Path, help="단일 PPTX의 출력 JSONL 경로")
    parser.add_argument("--output-dir", type=Path, default=Path("claims"), help="출력 폴더")
    parser.add_argument("--overwrite", action="store_true", help="기존 결과 덮어쓰기")
    parser.add_argument("--inspect", action="store_true", help="추출한 슬라이드 구조만 화면에 출력")
    return parser.parse_args(argv)


def main(argv: Sequence[str] | None = None) -> int:
    args = parse_arguments(argv)
    try:
        decks = find_inputs(args.inputs)
    except FileNotFoundError as exc:
        print(exc, file=sys.stderr)
        return 2
    if not decks:
        print("처리할 PPTX가 없습니다.", file=sys.stderr)
        return 2
    if args.output and len(decks) != 1:
        print("--output은 PPTX 한 개를 처리할 때만 사용할 수 있습니다.", file=sys.stderr)
        return 2

    failures = 0
    for deck in decks:
        output = args.output or args.output_dir / f"{deck.stem}.jsonl"
        if output.exists() and not args.overwrite and not args.inspect:
            print(f"SKIP {deck.name}: {output} 파일이 이미 있습니다.")
            continue
        try:
            slides = read_slides(deck)
            if args.inspect:
                print(json.dumps([asdict(slide) for slide in slides], ensure_ascii=False, indent=2))
                continue
            rows = [row for slide in slides for row in build_records(slide)]
            count = save_jsonl(output, rows)
            print(f"OK   {deck.name}: {count} claims -> {output}")
            for slide in slides:
                for warning in slide.warnings:
                    print(f"WARN {deck.name}: slide {slide.number}: {warning}", file=sys.stderr)
        except Exception as exc:  # 한 덱이 실패해도 나머지 덱은 계속 처리한다.
            failures += 1
            print(f"FAIL {deck.name}: {exc}", file=sys.stderr)
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
