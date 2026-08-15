"""manifest.csv 관리 도구 — 기계적으로 알 수 있는 값은 자동으로 채우고,
사람이 결정해야 하는 값은 빠뜨리지 않았는지 확인한다.

`docs/manifest.csv` 는 계획서(3장)가 "모든 작업의 기준 파일"이라고 못박은 파일이다.
그런데 지금까지는 새 덱/원문이 생길 때마다 사람이 CSV 를 손으로 열어 행을 추가해왔고,
등록을 건너뛰면 `build_queries.py` 가 doc_id 를 조용히 잘못 추정한 채로 그냥 진행해버린다
(`bio02_claudecode X1` 같은 경우가 실제로 이렇게 걸렸다). 이 스크립트는 그 틈을 좁힌다.

자동으로 채우는 것 — 파일을 열어서 세면 되는 값
  slide_count   decks/{deck_id}.pptx 를 열어 슬라이드 수를 센다
  pages         docs/raw/{doc_id}.pdf 를 열어 페이지 수를 센다
  lang          기본값 "ko" (이 프로젝트 문서는 전부 한국어)

절대 자동으로 채우지 않는 것 — 사람이 결정해야 하는 값
  doc_id        문서를 뭐라고 부를지는 팀 합의 사항이다. 파일명만 보고 추정하면
                `deck_id.split("__")[0]` 식 규칙이 깨지는 경우(공백·버전 접미사 등)
                엉뚱한 doc_id 를 만들어낸다.
  title/source/domain/owner/license/source_url
                원문에 대한 서지 정보 — PPT/PDF 안에 기계적으로 정확히 들어있지 않다.
                (license 는 특히 "저작권 이슈 없는 것만" 이라는 계획서 조건을 사람이
                직접 확인해야 하는 값이라 자동화 대상이 될 수 없다)
  split         train/test/generalization 배정은 8/11 회의에서 정한 편집 결정이다.
                자동 배정하면 "학습에 안 쓸 분야"가 실수로 학습에 섞일 수 있다.

사용법
  # 1. 등록 안 된 덱/원문이 있는지 훑어본다 (아무것도 고치지 않는다 — 매주 이거부터 돌려라)
  #    등록 안 된 덱마다 --add 명령을 통째로, 한 줄로 찍어준다. bash·PowerShell·cmd
  #    아무 데나 그대로 붙여넣고 FILL_ 로 시작하는 자리만 실제 값으로 바꾸면 된다.
  #    (한 줄인 이유: PowerShell 은 bash 처럼 줄 끝 `\` 로 이어붙이지 않는다 — 백틱을 쓴다.
  #    셸마다 문법이 달라서 아예 한 줄로 만들었다)
  python src/check_manifest.py --scan

  # 2. --scan 이 찍어준 명령을 그대로 붙여넣거나, 손으로 하나씩 추가한다. (한 줄)
  python src/check_manifest.py --add --deck-id "bio02_claudecode X1" --doc-id bio_02 --title "BF-7, 허혈성 뇌손상을 막고 기억력을 살리다" --source "대한해부학회지 38(2), 181-188, 2005" --domain 의학 --license KOGL-1 --tool claudecode --split train
"""
from __future__ import annotations

import argparse
import csv
import shlex
import sys
from pathlib import Path

FIELDS = ["doc_id", "title", "source", "domain", "owner", "license", "source_url",
          "lang", "pages", "split", "deck_id", "tool", "prompt_id", "slide_count"]

# 사람이 직접 결정해야 해서, 안 채우면 --add 실행 시 경고를 띄우는 필드
HUMAN_REQUIRED = ["title", "source", "domain", "license", "split"]


def read_manifest(path: Path) -> list[dict]:
    """utf-8-sig 로 연다 — 엑셀의 "CSV UTF-8" 저장이 파일 맨 앞에 붙이는 BOM을
    벗겨내기 위해서다. BOM 을 못 벗기면 첫 컬럼 헤더가 "doc_id" 가 아니라
    "\\ufeffdoc_id" 로 읽혀서, row.get("doc_id") 가 매 행 None 을 돌려주고
    doc_id 등록이 전부 안 된 것처럼 보인다(실제로 이 버그로 8개 덱이 다 깨졌었다).
    utf-8-sig 는 BOM 이 있으면 벗기고 없으면 그냥 utf-8 처럼 읽어서 어느 쪽이든 안전하다.
    """
    if not path.exists():
        return []
    with path.open(encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def write_manifest(path: Path, rows: list[dict]) -> None:
    """utf-8-sig 로 쓴다 — BOM 을 붙여야 엑셀이 더블클릭으로 열 때 한글을 안 깨뜨리고
    UTF-8 로 알아서 인식한다(BOM 없이 순수 UTF-8 로 쓰면 엑셀이 시스템 로캘(CP949)로
    잘못 해석해서 한글이 깨지는 사고가 난다 — 지난번 UnicodeDecodeError 가 그거였다).
    read_manifest() 가 BOM 유무 상관없이 읽으니 이후 어느 쪽으로도 문제없다.
    """
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        w = csv.DictWriter(f, fieldnames=FIELDS)
        w.writeheader()
        for r in rows:
            w.writerow({k: r.get(k, "") for k in FIELDS})


def slide_count_of(pptx_path: Path) -> int | None:
    try:
        from pptx import Presentation
        return len(Presentation(pptx_path).slides)
    except Exception as exc:  # 파일이 없거나 손상됐으면 사람이 보게 놔둔다
        print(f"     (slide_count 자동 계산 실패: {exc})", file=sys.stderr)
        return None


def page_count_of(pdf_path: Path) -> int | None:
    try:
        from pypdf import PdfReader
        return len(PdfReader(pdf_path).pages)
    except Exception as exc:
        print(f"     (pages 자동 계산 실패: {exc})", file=sys.stderr)
        return None


def guess_title(pptx_path: Path) -> str | None:
    """첫 슬라이드 제목을 읽어본다. --add 의 --title 초안으로만 쓴다 — 사람이 확인해야 한다."""
    try:
        from pptx import Presentation
        slides = Presentation(pptx_path).slides
        if not slides:
            return None
        slide = slides[0]
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.has_text_frame:
            text = " ".join(title_shape.text_frame.text.split())
            if text:
                return text
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = " ".join(shape.text_frame.text.split())
                if text:
                    return text
    except Exception:
        pass
    return None


def guess_tool(deck_id: str) -> str | None:
    """기존 규칙(`{doc_id}__{tool}`)을 따르는 deck_id 에서만 tool 을 뽑는다."""
    if "__" in deck_id:
        return deck_id.rsplit("__", 1)[1]
    return None


def guess_doc_id(deck_id: str) -> str | None:
    """기존 규칙을 따르는 deck_id 에서만 doc_id 를 뽑는다. 못 뽑으면 사람이 정한다."""
    if "__" in deck_id:
        return deck_id.split("__")[0]
    return None


# FILL_ 로 시작하는 토큰만 쓴다 — <...>, [...], $(...) 같은 건 bash 에선 괜찮아도
# PowerShell/cmd 에서는 리다이렉션·타입 캐스트·서브식으로 해석돼서 그대로 붙여넣으면
# 터미널이 깨진다(실제로 이 문제로 --add 가 실패했었다). FILL_ 토큰은 공백·특수문자가
# 없는 그냥 문자열이라, 안 고치고 그대로 실행해도 셸이 죽지 않고 manifest.csv 에
# "FILL_SOURCE" 같은 값이 그대로 들어갈 뿐이다 — 실행 전에 사람이 바꾸면 된다.
PLACEHOLDER = {
    "doc_id": "FILL_DOC_ID",
    "title": "FILL_TITLE",
    "source": "FILL_SOURCE",
    "domain": "FILL_DOMAIN",
    "license": "FILL_LICENSE",
    "tool": "FILL_TOOL",
    "prompt_id": "FILL_PROMPT_ID",
    "split": "FILL_SPLIT",
}

HINTS = {
    "source": "예: 학회지 45(1), 2025, 197-220",
    "domain": "예: 의학, 사회학, 예술",
    "license": "예: KOGL-1, CC-BY, CC BY-NC 4.0",
    "prompt_id": "예: P1",
    "split": "train / test / generalization 중 하나",
}


def add_command_for(deck_path: Path, doc_id_hint: str | None = None) -> str:
    """등록 안 된 덱 하나에 대해, 터미널에 바로 붙여넣을 수 있는 --add 명령을 한 줄로 만든다.

    한 줄로 만드는 이유: bash 는 줄 끝 `\\` 로 여러 줄을 이어 붙이지만 PowerShell 은
    그 문법이 없다(백틱을 쓴다). 셸을 안 가리려면 아예 한 줄이 제일 안전하다.
    기계적으로 뽑을 수 있는 값(deck_id, tool 추정, title 추정)은 채우고, 사람이 정해야
    하는 값은 FILL_ 로 시작하는 자리로 남긴다.
    """
    deck_id = deck_path.stem
    doc_id = doc_id_hint or guess_doc_id(deck_id) or PLACEHOLDER["doc_id"]
    title = guess_title(deck_path) or PLACEHOLDER["title"]
    tool = guess_tool(deck_id) or PLACEHOLDER["tool"]

    def val(v: str) -> str:
        return v if v.startswith("FILL_") else shlex.quote(v)

    parts = [
        "python src/check_manifest.py --add",
        f"--deck-id {shlex.quote(deck_id)}",
        f"--doc-id {val(doc_id)}",
        f"--title {val(title)}",
        f"--source {PLACEHOLDER['source']}",
        f"--domain {PLACEHOLDER['domain']}",
        f"--license {PLACEHOLDER['license']}",
        f"--tool {val(tool)}",
        f"--prompt-id {PLACEHOLDER['prompt_id']}",
        f"--split {PLACEHOLDER['split']}",
    ]
    return " ".join(parts)


# ─────────────────────────────────────────────────────────────
# --scan  : 등록 안 된 덱 · 원문을 찾아서 --add 명령 템플릿까지 찍어준다
# ─────────────────────────────────────────────────────────────

def scan(root: Path) -> int:
    rows = read_manifest(root / "docs" / "manifest.csv")
    known_decks = {r["deck_id"] for r in rows if r.get("deck_id")}
    known_docs = {r["doc_id"] for r in rows if r.get("doc_id")}

    decks_dir = root / "decks"
    raw_dir = root / "docs" / "raw"
    new_decks = sorted(p for p in decks_dir.glob("*.pptx")
                       if not p.name.startswith("~$") and p.stem not in known_decks)
    # 원문 PDF 는 파일명이 doc_id 와 다를 수 있다(bio02_text (예린bio).pdf 처럼).
    # 그래서 "매칭 안 됨"만 알려주고 어떤 doc_id 인지는 추정하지 않는다.
    unmatched_raw = sorted(p for p in raw_dir.glob("*.pdf") if p.stem not in known_docs)

    print(f"manifest.csv 등록된 덱 {len(known_decks)}개, 문서 {len(known_docs)}개\n")

    if not new_decks and not unmatched_raw:
        print("등록 안 된 덱/원문 없음.")
        return 0

    if new_decks:
        print(f"[등록 안 된 덱] decks/ 에는 있는데 manifest.csv deck_id 에는 없음 — {len(new_decks)}개")
        print("아래 명령은 한 줄이라 bash/PowerShell/cmd 어디에 붙여넣어도 그대로 동작한다.")
        print("FILL_ 로 시작하는 자리만 실제 값으로 바꿔라"
              "(deck_id·title·tool 은 추정해서 미리 채웠다 — 틀렸으면 고칠 것).")
        hint_line = "  ".join(f"{k}={v}" for k, v in HINTS.items())
        print(f"  참고: {hint_line}\n")
        for p in new_decks:
            n = slide_count_of(p)
            print(f"# \"{p.name}\"  (slide_count={n if n is not None else '?'} — 실행하면 자동으로 채워짐)")
            print(add_command_for(p))
            print()

    if unmatched_raw:
        print(f"[매칭 안 된 원문] docs/raw/ 파일명이 manifest.csv doc_id 어디와도 안 맞음 — {len(unmatched_raw)}개")
        print("  (어느 덱의 원문인지, doc_id 를 뭐로 할지는 파일명만으로 못 정한다. 사람이 확인 후")
        print("   위 --add 명령의 --doc-id 를 여기 맞춰 넣고, 파일도 docs/raw/{doc_id}.pdf 로 옮길 것)")
        for p in unmatched_raw:
            n = page_count_of(p)
            print(f"  - \"{p.name}\"  (pages={n if n is not None else '?'})")
        print()

    return 1


# ─────────────────────────────────────────────────────────────
# --add   : 행 하나를 추가한다
# ─────────────────────────────────────────────────────────────

def add(root: Path, args: argparse.Namespace) -> int:
    manifest_path = root / "docs" / "manifest.csv"
    rows = read_manifest(manifest_path)

    if any(r.get("deck_id") == args.deck_id for r in rows):
        print(f"FAIL {args.deck_id}: manifest.csv 에 이미 있다. 손으로 수정하거나 지우고 다시 추가할 것.",
              file=sys.stderr)
        return 2

    deck_path = root / "decks" / f"{args.deck_id}.pptx"
    slide_count = slide_count_of(deck_path) if deck_path.exists() else None
    if slide_count is None:
        print(f"WARN decks/{args.deck_id}.pptx 를 못 찾았다 — slide_count 는 비워둔다. "
              f"파일명이 정확한지 확인할 것(공백/버전 접미사 등).", file=sys.stderr)

    pdf_path = root / "docs" / "raw" / f"{args.doc_id}.pdf"
    pages = page_count_of(pdf_path) if pdf_path.exists() else None

    valid_splits = {"train", "test", "generalization"}
    split = args.split or ""
    if split and split not in valid_splits:
        print(f"WARN --split 값 '{split}' 은 train/test/generalization 중 하나가 아니다"
              f"(FILL_SPLIT 을 안 바꾸고 그대로 돌렸을 수 있다) — 빈 칸으로 저장한다.",
              file=sys.stderr)
        split = ""

    # FILL_ 로 시작하는 값은 안 바꾸고 그대로 돌린 것이다 — CSV 에는 글자 그대로 박아넣지
    #않고 빈 칸으로 남긴다. "FILL_SOURCE" 같은 쓰레기 값이 manifest.csv 에 남는 것보다
    # 빈 칸으로 두고 나중에 채워야 한다는 걸 명확히 하는 편이 낫다.
    def clean(v: str | None) -> str:
        v = v or ""
        return "" if v.startswith("FILL_") else v

    row = {
        "doc_id": args.doc_id,
        "title": clean(args.title),
        "source": clean(args.source),
        "domain": clean(args.domain),
        "owner": clean(args.owner),
        "license": clean(args.license),
        "source_url": clean(args.source_url),
        "lang": args.lang,
        "pages": pages if pages is not None else "",
        "split": split,
        "deck_id": args.deck_id,
        "tool": clean(args.tool),
        "prompt_id": clean(args.prompt_id),
        "slide_count": slide_count if slide_count is not None else "",
    }
    rows.append(row)
    write_manifest(manifest_path, rows)

    missing = [f for f in HUMAN_REQUIRED if not row[f]]
    print(f"OK   {args.deck_id} (doc_id={args.doc_id}) -> {manifest_path} 에 추가함")
    print(f"     slide_count={row['slide_count'] or '?'}  pages={row['pages'] or '?'}")
    if missing:
        print(f"\n!! 아직 안 채운 필수 항목: {', '.join(missing)}")
        print("   manifest.csv 를 열어 직접 채워라 — 자동으로 못 채우는 값이다(위 docstring 참고).")
        return 1
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--root", type=Path, default=Path("."))
    ap.add_argument("--scan", action="store_true", help="등록 안 된 덱/원문을 찾아서 보고만 한다")
    ap.add_argument("--add", action="store_true", help="새 행을 manifest.csv 에 추가한다")

    ap.add_argument("--deck-id", help="decks/{deck-id}.pptx 의 파일명(확장자 제외)")
    ap.add_argument("--doc-id", help="이 문서를 부를 이름 — 팀이 정한다")
    ap.add_argument("--title")
    ap.add_argument("--source", help="서지 정보, 예: '대한해부학회지 38(2), 181-188, 2005'")
    ap.add_argument("--domain")
    ap.add_argument("--owner")
    ap.add_argument("--license", help="예: KOGL-1, CC-BY")
    ap.add_argument("--source-url")
    ap.add_argument("--lang", default="ko")
    ap.add_argument("--split", help="train / test / generalization 중 하나 "
                                    "(choices 로 강제하지 않는다 — FILL_SPLIT 을 안 바꾸고 "
                                    "그대로 돌려도 크래시 대신 경고만 뜨게 하려고)")
    ap.add_argument("--tool", help="예: claudecode, chatGPT")
    ap.add_argument("--prompt-id")

    args = ap.parse_args()

    if args.scan:
        return scan(args.root)
    if args.add:
        if not args.deck_id or not args.doc_id:
            ap.error("--add 는 --deck-id 와 --doc-id 가 둘 다 필요하다")
        if args.doc_id.startswith("FILL_"):
            ap.error(f"--doc-id 가 아직 '{args.doc_id}' 그대로다 — 실제 doc_id 로 바꾸고 다시 실행할 것 "
                     f"(이 값은 passages/{{doc_id}}.jsonl 등 다른 단계에서도 그대로 쓰여서 "
                     f"placeholder 로 등록하면 안 된다)")
        return add(args.root, args)
    ap.error("--scan 또는 --add 가 필요하다")


if __name__ == "__main__":
    raise SystemExit(main())
